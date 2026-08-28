"""
DATA WHERE HOUSE? — PPTX Slide Deck Exporter
Converts the pitch slides into a 16:9 widescreen PowerPoint (.pptx) presentation.
"""

import os
import sys
import io
import asyncio
from pptx import Presentation
from pptx.util import Inches
import pypdfium2 as pdfium

# Import PDF export function from export_pdf
from export_pdf import export_pdf

async def generate_pptx(output_pptx_path=None, pdf_source_path=None):
    current_dir = os.path.dirname(os.path.abspath(__file__))
    
    if not pdf_source_path:
        pdf_source_path = os.path.join(current_dir, "DATA_WHERE_HOUSE_Pitch_Slides.pdf")
    
    if not output_pptx_path:
        output_pptx_path = os.path.join(current_dir, "DATA_WHERE_HOUSE_Pitch_Slides.pptx")

    print("Step 1: Generating/Updating master PDF from HTML slides...")
    await export_pdf(output_path=pdf_source_path)

    print(f"Step 2: Rendering high-resolution slides from PDF ({pdf_source_path})...")
    pdf = pdfium.PdfDocument(pdf_source_path)
    total_pages = len(pdf)
    print(f"Detected {total_pages} slides.")

    # Initialize PowerPoint presentation in 16:9 widescreen
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    for idx, page in enumerate(pdf):
        print(f"Processing slide {idx + 1}/{total_pages} at 4K resolution (scale=3)...")
        # Render at scale=3 for ultra-crisp display (approx 2880x1620)
        image = page.render(scale=3).to_pil()
        
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG', optimize=True)
        img_byte_arr.seek(0)

        # Add slide to presentation
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_byte_arr,
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    prs.save(output_pptx_path)
    file_size = os.path.getsize(output_pptx_path)
    print(f"\n[SUCCESS] Exported 16:9 PowerPoint presentation to:\n  -> {output_pptx_path} ({file_size:,} bytes)")
    return output_pptx_path

if __name__ == "__main__":
    asyncio.run(generate_pptx())
